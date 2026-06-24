#!/usr/bin/env python3
from __future__ import annotations

import argparse
import html
import json
import math
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from dataclasses import dataclass
from datetime import UTC, datetime
from pathlib import Path
from typing import Any
from urllib.parse import urljoin

import requests
from dotenv import load_dotenv


ROOT = Path(__file__).resolve().parents[1]
CODEX_TEMPLATE_DIR = ROOT / "codex"

QUANTITY_LABELS = {
    "pages": "صفحه",
    "words": "کلمه",
    "slides": "اسلاید",
}

ORDER_FIELDS = [
    ("id", "شناسه سفارش"),
    ("order_type", "نوع سفارش"),
    ("title", "عنوان"),
    ("title_english", "عنوان انگلیسی"),
    ("degree", "مقطع"),
    ("university", "دانشگاه"),
    ("field_of_study", "رشته یا گرایش"),
    ("faculty", "دانشکده"),
    ("department", "گروه آموزشی"),
    ("advisor_name", "استاد راهنما"),
    ("consultant_name", "استاد مشاور"),
    ("instructor_name", "استاد درس"),
    ("course_name", "نام درس"),
    ("language", "زبان"),
    ("academic_style", "شیوه ارجاع‌دهی"),
    ("methodology", "روش یا رویکرد"),
    ("image_count", "تعداد عکس موردنیاز"),
    ("deadline", "مهلت تحویل"),
    ("keywords", "کلیدواژه‌ها"),
]

REQUIRED_FIELDS_BY_ORDER_TYPE = {
    "پایان‌نامه کارشناسی": ["faculty", "advisor_name"],
    "پایان‌نامه کارشناسی ارشد": ["faculty", "advisor_name"],
    "رساله دکتری": ["faculty", "advisor_name"],
    "پروپوزال پایان‌نامه": ["faculty", "advisor_name", "abstract"],
    "تحقیق دانشگاهی": ["instructor_name", "course_name"],
    "ارائه و پاورپوینت": ["course_name"],
}


@dataclass(frozen=True)
class Config:
    backend_url: str
    worker_api_key: str
    worker_id: str
    workspace: Path

    @property
    def headers(self) -> dict[str, str]:
        return {"Authorization": f"Bearer {self.worker_api_key}"}


@dataclass(frozen=True)
class FigureAsset:
    rid: str
    target: str
    extension: str
    content_type: str
    data: bytes
    alt: str
    caption: str
    source_url: str | None = None
    license_label: str | None = None
    author: str | None = None


def load_config() -> Config:
    load_dotenv(ROOT / ".env")
    backend_url = os.getenv("BACKEND_URL") or os.getenv("APP_URL", "http://localhost:5173")
    return Config(
        backend_url=backend_url.rstrip("/"),
        worker_api_key=os.getenv("WORKER_API_KEY", "local-worker-dev-key"),
        worker_id=os.getenv("WORKER_ID", "local-pc-1"),
        workspace=ROOT / os.getenv("WORKER_WORKSPACE", "workspace"),
    )


def api_url(config: Config, path: str) -> str:
    return f"{config.backend_url}{path}"


def compact_value(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value.strip()
    return str(value).strip()


def markdown_escape(value: Any) -> str:
    text = compact_value(value)
    return text.replace("|", "\\|").replace("\n", "<br>")


def quantity_text(order: dict[str, Any]) -> str:
    value = order.get("quantity_value")
    quantity_type = order.get("quantity_type")
    label = QUANTITY_LABELS.get(str(quantity_type), compact_value(quantity_type) or "واحد")
    if value is None or value == "":
        return "مشخص نشده"
    return f"{value} {label}"


def primary_direction(order: dict[str, Any]) -> str:
    language = compact_value(order.get("language"))
    if language == "انگلیسی":
        return "ltr"
    return "rtl"


def missing_order_fields(order: dict[str, Any]) -> list[str]:
    required = list(REQUIRED_FIELDS_BY_ORDER_TYPE.get(compact_value(order.get("order_type")), []))
    if order.get("quantity_value") in (None, ""):
        required.append("quantity_value")
    missing = []
    for field in required:
        if not compact_value(order.get(field)):
            missing.append(field)
    return missing


def normalized_order_context(order: dict[str, Any]) -> dict[str, Any]:
    return {
        "id": order.get("id"),
        "order_type": order.get("order_type"),
        "title": order.get("title"),
        "title_english": order.get("title_english"),
        "degree": order.get("degree"),
        "university": order.get("university"),
        "field_of_study": order.get("field_of_study"),
        "faculty": order.get("faculty"),
        "department": order.get("department"),
        "advisor_name": order.get("advisor_name"),
        "consultant_name": order.get("consultant_name"),
        "instructor_name": order.get("instructor_name"),
        "course_name": order.get("course_name"),
        "language": order.get("language"),
        "academic_style": order.get("academic_style"),
        "methodology": order.get("methodology"),
        "quantity_type": order.get("quantity_type"),
        "quantity_value": order.get("quantity_value"),
        "quantity_text": quantity_text(order),
        "image_count": order.get("image_count"),
        "deadline": order.get("deadline"),
        "keywords": order.get("keywords"),
        "abstract": order.get("abstract"),
        "notes": order.get("notes"),
        "primary_direction": primary_direction(order),
        "font_guidance": {
            "persian": ["Vazirmatn", "IRANSans", "B Nazanin", "B Mitra", "Tahoma", "Arial"],
            "english": ["Times New Roman", "Calibri", "Arial"],
        },
        "missing_order_fields": missing_order_fields(order),
        "files": order.get("files", []),
        "references": order.get("references", []),
    }


def order_context_markdown(order: dict[str, Any]) -> str:
    context = normalized_order_context(order)
    lines = [
        "# Order Context",
        "",
        "This file is generated by the local worker from `customer_input.json`. Use it before planning, drafting, formatting, or reporting.",
        "",
        "## Core Details",
        "",
        "| Field | Value |",
        "| --- | --- |",
    ]
    for key, label in ORDER_FIELDS:
        value = order.get(key)
        if compact_value(value):
            lines.append(f"| {label} | {markdown_escape(value)} |")
    lines.append(f"| حجم موردنیاز | {markdown_escape(context['quantity_text'])} |")
    lines.append(f"| جهت اصلی خروجی | `{context['primary_direction']}` |")

    if compact_value(order.get("abstract")):
        lines.extend(["", "## Abstract Or Problem Statement", "", markdown_escape(order.get("abstract"))])

    if compact_value(order.get("notes")):
        lines.extend(["", "## Customer Notes", "", markdown_escape(order.get("notes"))])

    missing = context["missing_order_fields"]
    lines.extend(["", "## Missing Or Unclear Order Details", ""])
    if missing:
        for field in missing:
            lines.append(f"- `{field}`")
    else:
        lines.append("- None detected by the worker precheck.")

    lines.extend(
        [
            "",
            "## Files",
            "",
            "| Type | Name | URL/Path |",
            "| --- | --- | --- |",
        ]
    )
    files = order.get("files", [])
    if files:
        for file_info in files:
            lines.append(
                "| "
                f"{markdown_escape(file_info.get('file_type'))} | "
                f"{markdown_escape(file_info.get('original_name'))} | "
                f"{markdown_escape(file_info.get('url'))} |"
            )
    else:
        lines.append("| - | No files listed | - |")

    lines.extend(
        [
            "",
            "## Customer References",
            "",
            "| Type | Title | Authors | Year | Required |",
            "| --- | --- | --- | --- | --- |",
        ]
    )
    references = order.get("references", [])
    if references:
        for reference in references:
            lines.append(
                "| "
                f"{markdown_escape(reference.get('reference_type'))} | "
                f"{markdown_escape(reference.get('title'))} | "
                f"{markdown_escape(reference.get('authors'))} | "
                f"{markdown_escape(reference.get('year'))} | "
                f"{'yes' if reference.get('required_usage') else 'no'} |"
            )
    else:
        lines.append("| - | No customer references listed | - | - | - |")

    lines.extend(
        [
            "",
            "## Output Style Requirements",
            "",
            "- Persian and Arabic content must be RTL.",
            "- English titles, URLs, DOIs, emails, code, file paths, citation keys, and IDs must remain LTR.",
            "- Preferred Persian fonts: Vazirmatn, IRANSans, B Nazanin, B Mitra, then Tahoma/Arial fallback.",
            "- Preferred English fonts: Times New Roman, Calibri, Arial.",
            "- Use clear headings, metadata tables, review checklists, and print-friendly spacing.",
            "- If a conversion tool cannot enforce fonts or directionality, record that limitation in `reports/compliance_report.md`.",
            "",
        ]
    )
    return "\n".join(lines)


def write_order_context(workspace: Path, order: dict[str, Any]) -> None:
    (workspace / "input").mkdir(parents=True, exist_ok=True)
    (workspace / "extracted").mkdir(parents=True, exist_ok=True)
    (workspace / "input" / "order_context.md").write_text(
        order_context_markdown(order),
        encoding="utf-8",
    )
    (workspace / "extracted" / "order_context.json").write_text(
        json.dumps(normalized_order_context(order), ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


def require_success(response: requests.Response) -> dict[str, Any]:
    if response.ok:
        if not response.content:
            return {}
        return response.json()

    detail: Any
    try:
        detail = response.json().get("detail", response.text)
    except ValueError:
        detail = response.text
    raise SystemExit(f"API error {response.status_code}: {detail}")


def order_dir(config: Config, order_id: str) -> Path:
    return config.workspace / f"order_{order_id}"


def normalize_workspace(workspace: str | None) -> Path | None:
    if not workspace:
        return None
    path = Path(workspace).expanduser()
    if not path.is_absolute():
        path = (Path.cwd() / path).resolve()
    return path


def workspace_from_cwd() -> Path | None:
    current = Path.cwd().resolve()
    if (current / "customer_input.json").exists():
        return current
    return None


def infer_order_id_from_workspace(workspace: Path) -> str | None:
    customer_input = workspace / "customer_input.json"
    if customer_input.exists():
        data = json.loads(customer_input.read_text(encoding="utf-8"))
        if data.get("id"):
            return str(data["id"])

    if workspace.name.startswith("order_"):
        return workspace.name.removeprefix("order_")

    return None


def resolve_order_context(
    config: Config,
    provided_order_id: str | None = None,
    provided_workspace: str | None = None,
) -> tuple[str, Path]:
    workspace = normalize_workspace(provided_workspace) or workspace_from_cwd()

    if provided_order_id:
        return provided_order_id, workspace or order_dir(config, provided_order_id)

    if workspace:
        inferred_order_id = infer_order_id_from_workspace(workspace)
        if inferred_order_id:
            return inferred_order_id, workspace

    raise SystemExit(
        "Could not infer the order. Run this command inside `worker/workspace/order_<id>`, "
        "or pass --workspace/--order-id explicitly."
    )


def ensure_workspace(config: Config, order: dict[str, Any]) -> Path:
    order_id = order["id"]
    workspace = order_dir(config, order_id)
    for relative in [
        "input/files",
        "input/references",
        "extracted",
        "planning",
        "drafts",
        "reports",
        "final",
    ]:
        (workspace / relative).mkdir(parents=True, exist_ok=True)

    (workspace / "customer_input.json").write_text(
        json.dumps(order, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    write_order_context(workspace, order)

    target_codex_dir = workspace / "codex"
    if target_codex_dir.exists():
        shutil.rmtree(target_codex_dir)
    shutil.copytree(CODEX_TEMPLATE_DIR, target_codex_dir)
    (workspace / "AGENTS.md").write_text(
        "# Order Workspace Instructions\n\n"
        "Read and follow `codex/AGENTS.md`, then execute "
        "`codex/workflows/order_workflow.md` through the human review package stage.\n",
        encoding="utf-8",
    )

    checklist = workspace / "reports" / "human_review_checklist.md"
    if not checklist.exists():
        checklist.write_text(
            "# Human Review Checklist\n\n## Missing or unclear inputs\n\n- Pending intake review.\n",
            encoding="utf-8",
        )

    return workspace


def download_file(config: Config, workspace: Path, file_info: dict[str, Any]) -> Path:
    url = urljoin(config.backend_url, file_info["url"])
    original_name = file_info["original_name"]
    target = workspace / "input" / "files" / f"{file_info['file_type']}__{original_name}"
    response = requests.get(url, timeout=120)
    response.raise_for_status()
    target.write_bytes(response.content)
    return target


def claim(config: Config) -> None:
    response = requests.post(
        api_url(config, "/api/worker/orders/claim-oldest"),
        headers=config.headers,
        json={"workerId": config.worker_id},
        timeout=30,
    )
    payload = require_success(response)
    order = payload["customerInput"]
    workspace = ensure_workspace(config, order)

    start_response = requests.post(
        api_url(config, f"/api/worker/orders/{order['id']}/start"),
        headers={**config.headers, "Content-Type": "application/json"},
        json={"workerId": config.worker_id, "notes": "Local worker workspace created."},
        timeout=30,
    )
    require_success(start_response)

    downloaded = []
    for file_info in order.get("files", []):
        downloaded.append(download_file(config, workspace, file_info))

    print(json.dumps(
        {
            "orderId": order["id"],
            "workspace": str(workspace),
            "downloadedFiles": [str(path) for path in downloaded],
            "nextStep": "cd into the workspace and run Codex with AGENTS.md",
        },
        ensure_ascii=False,
        indent=2,
    ))


def heartbeat(config: Config, args: argparse.Namespace) -> None:
    order_id, _workspace = resolve_order_context(config, args.order_id, args.workspace)
    response = requests.post(
        api_url(config, f"/api/worker/orders/{order_id}/heartbeat"),
        headers={**config.headers, "Content-Type": "application/json"},
        json={"workerId": config.worker_id},
        timeout=30,
    )
    payload = require_success(response)
    print(json.dumps({"orderId": order_id, "status": payload["status"]}, indent=2))


def submit_draft(config: Config, args: argparse.Namespace) -> None:
    order_id, workspace = resolve_order_context(config, args.order_id, args.workspace)
    draft_path = args.draft_file
    default_draft = workspace / "drafts" / "assisted_draft.md"
    if not draft_path and default_draft.exists():
        draft_path = str(default_draft)

    files = {}
    if draft_path:
        files["draft_file"] = open(draft_path, "rb")
    try:
        response = requests.post(
            api_url(config, f"/api/worker/orders/{order_id}/submit-draft"),
            headers=config.headers,
            data={"worker_id": config.worker_id, "notes": args.notes or ""},
            files=files,
            timeout=120,
        )
        payload = require_success(response)
        print(json.dumps({"orderId": order_id, "status": payload["status"]}, indent=2))
    finally:
        for handle in files.values():
            handle.close()


def resolved_upload(user_path: str | None, default_path: Path) -> str | None:
    if user_path:
        return user_path
    if default_path.exists():
        return str(default_path)
    return None


def optional_upload(files: dict[str, Any], field_name: str, path: str | None) -> None:
    if path:
        files[field_name] = open(path, "rb")


def default_source_path(workspace: Path) -> str | None:
    return resolved_upload(None, workspace / "final" / "deliverable_source.md") or resolved_upload(
        None, workspace / "final" / "thesis_source.md"
    )


def default_docx_path(workspace: Path) -> str | None:
    return resolved_upload(None, workspace / "final" / "deliverable.docx") or resolved_upload(
        None, workspace / "final" / "thesis.docx"
    )


def default_pdf_path(workspace: Path) -> str | None:
    return resolved_upload(None, workspace / "final" / "deliverable.pdf") or resolved_upload(
        None, workspace / "final" / "thesis.pdf"
    )


def workspace_order(workspace: Path) -> dict[str, Any]:
    customer_input = workspace / "customer_input.json"
    if not customer_input.exists():
        return {}
    return json.loads(customer_input.read_text(encoding="utf-8"))


def refresh_workspace_order_snapshot(workspace: Path, order: dict[str, Any]) -> None:
    if not order:
        return
    (workspace / "customer_input.json").write_text(
        json.dumps(order, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    write_order_context(workspace, order)


def is_presentation_order(order: dict[str, Any]) -> bool:
    order_type = compact_value(order.get("order_type"))
    quantity_type = compact_value(order.get("quantity_type"))
    return "پاورپوینت" in order_type or "ارائه" in order_type or quantity_type == "slides"


def validate_review_package(workspace: Path, order: dict[str, Any], upload_paths: dict[str, str | None]) -> None:
    required: list[tuple[str, str | None]] = [
        ("deliverable source", upload_paths.get("deliverable_source")),
        ("editable Word file", upload_paths.get("docx")),
        ("compliance report", upload_paths.get("compliance_report")),
        ("reference usage report", upload_paths.get("reference_usage_report")),
        ("human review checklist", upload_paths.get("human_review_checklist")),
        ("final README", upload_paths.get("final_readme")),
    ]
    if is_presentation_order(order):
        required.append(("editable PowerPoint file", upload_paths.get("pptx")))

    expected_figures = requested_or_default_image_count(order)
    if expected_figures > 0:
        required.append(("figure source metadata", upload_paths.get("image_sources")))

    missing = [label for label, path in required if not path or not Path(path).exists() or Path(path).stat().st_size == 0]
    if missing:
        bullet_list = "\n".join(f"- {item}" for item in missing)
        raise SystemExit(
            "Review package is incomplete; refusing to submit final package.\n"
            f"{bullet_list}\n"
            f"Workspace: {workspace}"
        )

    if expected_figures > 0:
        validate_image_sources_file(workspace, upload_paths.get("image_sources"), expected_figures)
    validate_docx_rtl_quality(order, upload_paths.get("docx"))
    validate_no_placeholder_outputs(upload_paths)
    if is_presentation_order(order):
        validate_pptx_ui_quality(order, upload_paths.get("pptx"), workspace)


def validate_image_sources_file(workspace: Path, image_sources_path: str | None, expected_count: int) -> None:
    if not image_sources_path:
        raise SystemExit("Figure source metadata is required when the order expects visuals.")

    path = Path(image_sources_path)
    try:
        sources = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, ValueError) as exc:
        raise SystemExit(f"Could not read figure source metadata: {path}") from exc

    actual_count = len(sources) if isinstance(sources, list) else 0
    if not isinstance(sources, list) or actual_count < expected_count:
        raise SystemExit(
            "Figure source metadata is incomplete; refusing to submit final package.\n"
            f"Expected at least {expected_count} sourced figure records, found {actual_count}."
        )

    problems: list[str] = []
    for index, source in enumerate(sources, start=1):
        if not isinstance(source, dict):
            problems.append(f"record {index}: not an object")
            continue

        local_path = compact_value(source.get("local_path"))
        source_url = compact_value(source.get("description_url")) or compact_value(source.get("url"))
        license_label = compact_value(source.get("license"))
        generated_flag = compact_value(source.get("generated")) or compact_value(source.get("generation_tool"))
        if not local_path:
            problems.append(f"record {index}: missing local_path")
        elif not (workspace / local_path).exists():
            problems.append(f"record {index}: local image file does not exist")
        elif "cycle" in local_path:
            try:
                from PIL import Image
                with Image.open(workspace / local_path) as image:
                    width, height = image.size
                if height > width:
                    problems.append(f"record {index}: cycle chart is portrait and will be hard to read on presentation slides")
            except Exception:
                problems.append(f"record {index}: could not inspect cycle chart dimensions")
        if not source_url:
            problems.append(f"record {index}: missing internet source URL")
        if not license_label:
            problems.append(f"record {index}: missing license label")
        if generated_flag:
            problems.append(f"record {index}: generated diagrams/images are not accepted for automatic final packaging")

    if problems:
        bullet_list = "\n".join(f"- {problem}" for problem in problems)
        raise SystemExit(
            "Figure source metadata failed validation; refusing to submit final package.\n"
            f"{bullet_list}"
        )


def validate_docx_rtl_quality(order: dict[str, Any], docx_path: str | None) -> None:
    if not docx_path or primary_direction(order) != "rtl":
        return

    path = Path(docx_path)
    if not path.exists():
        return

    try:
        with zipfile.ZipFile(path) as docx:
            document_xml = docx.read("word/document.xml").decode("utf-8")
            settings_xml = docx.read("word/settings.xml").decode("utf-8") if "word/settings.xml" in docx.namelist() else ""
            styles_xml = docx.read("word/styles.xml").decode("utf-8") if "word/styles.xml" in docx.namelist() else ""
    except (OSError, KeyError, zipfile.BadZipFile, UnicodeDecodeError) as exc:
        raise SystemExit(f"Could not inspect DOCX RTL formatting: {path}") from exc

    paragraph_count = max(document_xml.count("<w:p>"), document_xml.count("<w:p "))
    bidi_count = document_xml.count("<w:bidi") + styles_xml.count("<w:bidi")
    right_count = document_xml.count('w:jc w:val="right"') + styles_xml.count('w:jc w:val="right"')
    left_count = document_xml.count('w:jc w:val="left"')
    has_fa_lang = 'w:bidi="fa-IR"' in settings_xml + styles_xml + document_xml
    has_cs_font = 'w:cs="B Nazanin"' in styles_xml + document_xml or 'w:cs="Vazirmatn"' in styles_xml + document_xml

    problems: list[str] = []
    if paragraph_count and bidi_count < max(3, paragraph_count // 2):
        problems.append("most paragraphs are not marked as RTL/bidi")
    if paragraph_count and right_count < max(3, paragraph_count // 2):
        problems.append("most paragraphs are not explicitly right-aligned")
    if left_count:
        problems.append("left-aligned paragraph formatting was found in an RTL deliverable")
    if not has_fa_lang:
        problems.append("Persian complex-script language metadata is missing")
    if not has_cs_font:
        problems.append("Persian complex-script font metadata is missing")

    if problems:
        bullet_list = "\n".join(f"- {problem}" for problem in problems)
        raise SystemExit(
            "DOCX RTL/right-alignment validation failed; refusing to submit final package.\n"
            f"{bullet_list}"
        )


def docx_text(path: Path) -> str:
    try:
        with zipfile.ZipFile(path) as docx:
            document_xml = docx.read("word/document.xml").decode("utf-8")
    except (OSError, KeyError, zipfile.BadZipFile, UnicodeDecodeError):
        return ""
    return re.sub(r"<[^>]+>", " ", document_xml)


def pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    try:
        prs = Presentation(str(path))
    except Exception:
        return ""
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                lines.append(shape.text)
    return "\n".join(lines)


def text_for_output(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return docx_text(path)
    if suffix == ".pptx":
        return pptx_text(path)
    if suffix in {".md", ".txt", ".json"}:
        try:
            return path.read_text(encoding="utf-8", errors="ignore")
        except OSError:
            return ""
    return ""


def validate_no_placeholder_outputs(upload_paths: dict[str, str | None]) -> None:
    problems: list[str] = []
    for label, raw_path in upload_paths.items():
        if not raw_path:
            continue
        path = Path(raw_path)
        if not path.exists():
            continue
        text = text_for_output(path)
        for placeholder in FORBIDDEN_PLACEHOLDERS:
            if placeholder in text:
                problems.append(f"{label}: contains unfinished placeholder `{placeholder}`")
    if problems:
        bullet_list = "\n".join(f"- {problem}" for problem in problems)
        raise SystemExit(
            "Output placeholder validation failed; refusing to submit final package.\n"
            f"{bullet_list}"
        )


def validate_pptx_ui_quality(order: dict[str, Any], pptx_path: str | None, workspace: Path | None = None) -> None:
    if not pptx_path:
        return
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise SystemExit("python-pptx is required to validate presentation UI quality.") from exc

    path = Path(pptx_path)
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise SystemExit(f"Could not inspect PPTX UI quality: {path}") from exc

    problems: list[str] = []
    try:
        expected_slides = int(order.get("quantity_value") or 0) if compact_value(order.get("quantity_type")) == "slides" else 0
    except (TypeError, ValueError):
        expected_slides = 0
    if expected_slides and len(prs.slides) != expected_slides:
        problems.append(f"expected {expected_slides} slides, found {len(prs.slides)}")

    picture_count = 0
    for slide_index, slide in enumerate(prs.slides, start=1):
        text = "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))
        if not text.strip():
            problems.append(f"slide {slide_index}: no readable text")
        for placeholder in FORBIDDEN_PLACEHOLDERS:
            if placeholder in text:
                problems.append(f"slide {slide_index}: contains placeholder `{placeholder}`")
        picture_count += sum(1 for shape in slide.shapes if str(shape.shape_type) == "PICTURE (13)")
        if len(slide.shapes) < 4:
            problems.append(f"slide {slide_index}: too few layout elements; looks like raw content")
        if slide_index == 1:
            first_slide_pictures = sum(1 for shape in slide.shapes if str(shape.shape_type) == "PICTURE (13)")
            if first_slide_pictures == 0:
                problems.append("slide 1: title slide has no visual background or image")

    expected_figures = requested_or_default_image_count(order)
    if expected_figures and picture_count < min(expected_figures, 3):
        problems.append(f"expected at least {min(expected_figures, 3)} sourced visuals in PPTX, found {picture_count}")

    if workspace:
        sources_path = workspace / "final" / "figures" / "image_sources.json"
        if expected_figures and not sources_path.exists():
            problems.append("missing image_sources.json for presentation visuals")

    if problems:
        bullet_list = "\n".join(f"- {problem}" for problem in problems)
        raise SystemExit(
            "PPTX UI quality validation failed; refusing to submit final package.\n"
            f"{bullet_list}"
        )


def submit_final(config: Config, args: argparse.Namespace) -> None:
    order_id, workspace = resolve_order_context(config, args.order_id, args.workspace)
    order = workspace_order(workspace)
    deliverable_source = args.deliverable_source or default_source_path(workspace)
    pptx = resolved_upload(args.pptx, workspace / "final" / "deliverable.pptx")
    docx = args.docx or default_docx_path(workspace)
    pdf = args.pdf or default_pdf_path(workspace)
    compliance_report = resolved_upload(
        args.compliance_report, workspace / "reports" / "compliance_report.md"
    )
    reference_usage_report = resolved_upload(
        args.reference_usage_report, workspace / "reports" / "reference_usage_report.md"
    )
    human_review_checklist = resolved_upload(
        args.human_review_checklist, workspace / "reports" / "human_review_checklist.md"
    )
    final_readme = resolved_upload(args.final_readme, workspace / "final" / "README.md")
    image_sources = resolved_upload(
        args.image_sources, workspace / "final" / "figures" / "image_sources.json"
    )

    upload_paths = {
        "deliverable_source": deliverable_source,
        "pptx": pptx,
        "docx": docx,
        "pdf": pdf,
        "compliance_report": compliance_report,
        "reference_usage_report": reference_usage_report,
        "human_review_checklist": human_review_checklist,
        "final_readme": final_readme,
        "image_sources": image_sources,
    }
    if not args.skip_package_check:
        validate_review_package(workspace, order, upload_paths)

    files: dict[str, Any] = {}
    optional_upload(files, "deliverable_source", deliverable_source)
    optional_upload(files, "pptx_file", pptx)
    optional_upload(files, "docx_file", docx)
    optional_upload(files, "pdf_file", pdf)
    optional_upload(files, "compliance_report", compliance_report)
    optional_upload(files, "reference_usage_report", reference_usage_report)
    optional_upload(files, "human_review_checklist", human_review_checklist)
    optional_upload(files, "final_readme", final_readme)
    optional_upload(files, "image_sources", image_sources)
    if not files:
        raise SystemExit(
            "At least one output file is required. Expected defaults under "
            f"{workspace}/final/ and {workspace}/reports/, or pass file paths explicitly."
        )

    try:
        response = requests.post(
            api_url(config, f"/api/worker/orders/{order_id}/submit-final"),
            headers=config.headers,
            data={
                "worker_id": config.worker_id,
                "notes": args.notes or "",
                "replace_existing": "true" if args.replace_existing else "false",
            },
            files=files,
            timeout=180,
        )
        payload = require_success(response)
        refresh_workspace_order_snapshot(workspace, payload)
        print(json.dumps(
            {
                "orderId": order_id,
                "status": payload["status"],
                "finalOutputs": len(payload.get("final_outputs", [])),
            },
            indent=2,
        ))
    finally:
        for handle in files.values():
            handle.close()


def fail(config: Config, args: argparse.Namespace) -> None:
    order_id, _workspace = resolve_order_context(config, args.order_id, args.workspace)
    response = requests.post(
        api_url(config, f"/api/worker/orders/{order_id}/fail"),
        headers={**config.headers, "Content-Type": "application/json"},
        json={"workerId": config.worker_id, "notes": args.notes or "Local worker marked failed."},
        timeout=30,
    )
    payload = require_success(response)
    print(json.dumps({"orderId": order_id, "status": payload["status"]}, indent=2))


def reset(config: Config, args: argparse.Namespace) -> None:
    order_id, _workspace = resolve_order_context(config, args.order_id, args.workspace)
    response = requests.post(
        api_url(config, f"/api/worker/orders/{order_id}/reset"),
        headers={**config.headers, "Content-Type": "application/json"},
        json={"workerId": config.worker_id, "notes": args.notes or "Local worker reset order."},
        timeout=30,
    )
    payload = require_success(response)
    print(json.dumps({"orderId": order_id, "status": payload["status"]}, indent=2))


def contains_rtl(text: str) -> bool:
    return any("\u0600" <= char <= "\u06ff" or "\u0750" <= char <= "\u077f" for char in text)


def docx_paragraph(line: str, heading: bool = False) -> str:
    rtl = contains_rtl(line)
    alignment = "right" if rtl else "left"
    bidi = "<w:bidi/>" if rtl else ""
    size = "32" if heading else "24"
    bold = "<w:b/>" if heading else ""
    paragraph_props = f'<w:pPr>{bidi}<w:jc w:val="{alignment}"/></w:pPr>'
    run_props = (
        '<w:rPr>'
        '<w:rFonts w:ascii="Times New Roman" w:hAnsi="Times New Roman" w:cs="B Nazanin"/>'
        f"{bold}<w:sz w:val=\"{size}\"/><w:szCs w:val=\"{size}\"/>"
        '</w:rPr>'
    )
    return (
        f"<w:p>{paragraph_props}<w:r>{run_props}<w:t xml:space=\"preserve\">"
        + html.escape(line)
        + "</w:t></w:r></w:p>"
    )


def write_minimal_docx(path: Path, title: str, body_lines: list[str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    paragraphs = [docx_paragraph(title, heading=True)]
    paragraphs.extend(docx_paragraph(line) for line in body_lines)
    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        "<w:body>"
        + "".join(paragraphs)
        + '<w:sectPr><w:pgSz w:w="11906" w:h="16838"/><w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440"/></w:sectPr>'
        + "</w:body></w:document>"
    )
    content_types = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        "</Types>"
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
        'Target="word/document.xml"/>'
        "</Relationships>"
    )
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as docx:
        docx.writestr("[Content_Types].xml", content_types)
        docx.writestr("_rels/.rels", rels)
        docx.writestr("word/document.xml", document_xml)


def write_minimal_pdf(path: Path, title: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    safe_title = title.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    pdf = (
        b"%PDF-1.4\n"
        b"1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj\n"
        b"2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj\n"
        b"3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >> endobj\n"
        + f"4 0 obj << /Length {len(safe_title) + 76} >> stream\n".encode()
        + f"BT /F1 18 Tf 72 720 Td ({safe_title}) Tj 0 -28 Td "
        "(Mock academic package for system testing only.) Tj ET\n".encode()
        + b"endstream endobj\n"
        b"5 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj\n"
        b"xref\n0 6\n0000000000 65535 f \n"
        b"trailer << /Root 1 0 R /Size 6 >>\nstartxref\n0\n%%EOF\n"
    )
    path.write_bytes(pdf)


def mock_generate(config: Config, args: argparse.Namespace) -> None:
    order_id, workspace = resolve_order_context(config, args.order_id, args.workspace)
    customer_input = workspace / "customer_input.json"
    order = json.loads(customer_input.read_text(encoding="utf-8")) if customer_input.exists() else {}
    title = order.get("title") or "Mock Academic Deliverable"
    generated_at = datetime.now(UTC).isoformat()

    (workspace / "extracted").mkdir(parents=True, exist_ok=True)
    (workspace / "planning").mkdir(parents=True, exist_ok=True)
    (workspace / "drafts").mkdir(parents=True, exist_ok=True)
    (workspace / "reports").mkdir(parents=True, exist_ok=True)
    (workspace / "final").mkdir(parents=True, exist_ok=True)
    write_order_context(workspace, order)

    (workspace / "extracted" / "university_rules.json").write_text(
        json.dumps(
            {
                "source_files": [],
                "order_context_used": normalized_order_context(order),
                "degree": order.get("degree"),
                "university": order.get("university"),
                "language": order.get("language"),
                "citation_style": order.get("academic_style"),
                "uncertain_rules": ["Mock test package; no real guideline extraction performed."],
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    (workspace / "extracted" / "references.json").write_text(
        json.dumps(
            {
                "references": order.get("references", []),
                "missing_or_unreadable_sources": [],
                "mock": True,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    (workspace / "planning" / "outline.md").write_text(
        f"# {title}\n\n"
        "## Order Metadata\n\n"
        f"- Type: {order.get('order_type') or '-'}\n"
        f"- University: {order.get('university') or '-'}\n"
        f"- Field: {order.get('field_of_study') or '-'}\n"
        f"- Quantity: {quantity_text(order)}\n\n"
        "## Mock Structure\n\n"
        "1. Introduction\n2. Literature Snapshot\n3. Mock Methodology\n4. Mock Findings\n5. Conclusion\n",
        encoding="utf-8",
    )
    (workspace / "planning" / "chapter_plan.md").write_text(
        "# Chapter Plan\n\n"
        f"- Requested quantity: {quantity_text(order)}\n"
        "- This is a tiny mock chapter plan for system testing only.\n",
        encoding="utf-8",
    )
    (workspace / "drafts" / "assisted_draft.md").write_text(
        f"# {title}\n\n"
        "| Field | Value |\n"
        "| --- | --- |\n"
        f"| Order type | {order.get('order_type') or '-'} |\n"
        f"| University | {order.get('university') or '-'} |\n"
        f"| Advisor/Instructor | {order.get('advisor_name') or order.get('instructor_name') or '-'} |\n"
        f"| Quantity | {quantity_text(order)} |\n\n"
        "This is a simple mock academic deliverable generated only to test the order system.\n\n"
        "It is not real academic work and must not be submitted as a finished assignment.\n",
        encoding="utf-8",
    )
    (workspace / "reports" / "compliance_report.md").write_text(
        f"# Compliance Report\n\nGenerated: {generated_at}\n\n"
        "- Mock package created successfully.\n"
        "- Real compliance review was not performed.\n",
        encoding="utf-8",
    )
    (workspace / "reports" / "reference_usage_report.md").write_text(
        f"# Reference Usage Report\n\nGenerated: {generated_at}\n\n"
        "- Mock package only; no real reference usage analysis performed.\n",
        encoding="utf-8",
    )
    (workspace / "reports" / "human_review_checklist.md").write_text(
        "# Human Review Checklist\n\n- This is a mock package for system testing only.\n",
        encoding="utf-8",
    )
    (workspace / "final" / "README.md").write_text(
        f"# Final Package\n\n"
        f"Order: {order_id}\n\n"
        f"Type: {order.get('order_type') or '-'}\n\n"
        f"Title: {title}\n\n"
        f"Quantity: {quantity_text(order)}\n\n"
        "Mock final files generated for system testing. Human review is still required before upload.\n",
        encoding="utf-8",
    )
    write_minimal_docx(
        workspace / "final" / "deliverable.docx",
        title,
        [
            "Mock academic deliverable package for system testing only.",
            "This file verifies final output upload and download.",
        ],
    )
    write_minimal_pdf(workspace / "final" / "deliverable.pdf", title)

    print(
        json.dumps(
            {
                "orderId": order_id,
                "workspace": str(workspace),
                "generated": [
                    "final/deliverable.docx",
                    "final/deliverable.pdf",
                    "reports/compliance_report.md",
                    "reports/reference_usage_report.md",
                ],
            },
            ensure_ascii=False,
            indent=2,
        )
    )


def markdown_to_plain_lines(markdown: str) -> list[str]:
    lines: list[str] = []
    for raw_line in markdown.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        line = line.lstrip("#").strip()
        line = line.replace("**", "").replace("__", "").replace("`", "")
        if line.startswith("- "):
            line = line[2:].strip()
        lines.append(line)
    return lines or ["Academic source was empty."]


def requested_page_count(order: dict[str, Any]) -> int | None:
    if order.get("quantity_type") != "pages":
        return None
    try:
        value = int(order.get("quantity_value"))
    except (TypeError, ValueError):
        return None
    return value if value > 0 else None


IMAGE_CONTENT_TYPES = {
    ".svg": "image/svg+xml",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
}

COMMONS_API_URL = "https://commons.wikimedia.org/w/api.php"
COMMONS_USER_AGENT = "payanname-worker/1.0 (academic document image research)"
FORBIDDEN_PLACEHOLDERS = [
    "برای مرور انسانی تکمیل شود",
    "تکمیل شود",
    "TODO",
    "TBD",
    "[NEEDS",
    "در دسترس نیست",
]


def requested_or_default_image_count(order: dict[str, Any]) -> int:
    explicit = order.get("image_count")
    if explicit is not None and explicit != "":
        try:
            return max(int(explicit), 0)
        except (TypeError, ValueError):
            return 0

    quantity_type = compact_value(order.get("quantity_type"))
    try:
        quantity_value = int(order.get("quantity_value") or 0)
    except (TypeError, ValueError):
        quantity_value = 0

    order_type = compact_value(order.get("order_type"))
    if quantity_type == "slides" and quantity_value:
        return min(max(math.ceil(quantity_value / 4), 1), 6)
    if quantity_type == "pages" and quantity_value:
        return min(max(math.ceil(quantity_value / 3), 1), 4)
    if quantity_type == "words" and quantity_value:
        return min(max(math.ceil(quantity_value / 1200), 1), 4)
    if "پایان" in order_type or "رساله" in order_type or "پروپوزال" in order_type:
        return 3
    return 1


def strip_html(value: Any) -> str:
    text = compact_value(value)
    text = re.sub(r"<[^>]+>", "", text)
    return html.unescape(text).strip()


def image_search_queries(order: dict[str, Any], headings: list[str]) -> list[str]:
    title = compact_value(order.get("title"))
    title_english = compact_value(order.get("title_english"))
    keywords = compact_value(order.get("keywords"))
    queries: list[str] = []
    combined = " ".join([title, title_english, keywords, *headings[:2]])
    if any(term in combined for term in ["قالب تونلی", "تونلی", "Tunnel", "tunnel"]):
        queries.extend([
            "tunnel formwork concrete construction",
            "concrete formwork rapid housing construction",
            "tunnel form construction concrete wall slab",
        ])
    if any(term in combined for term in ["قالب", "formwork", "Formwork"]):
        queries.extend([
            "concrete formwork construction",
            "formwork reinforced concrete construction",
        ])

    for value in [title_english, title, keywords]:
        if value:
            queries.append(value)

    normalized: list[str] = []
    seen: set[str] = set()
    for query in queries:
        query = re.sub(r"\s+", " ", query).strip()
        if query and query.casefold() not in seen:
            seen.add(query.casefold())
            normalized.append(query)
    return normalized[:8]


def relevant_commons_candidate(candidate: dict[str, Any], query: str) -> bool:
    text = " ".join([
        compact_value(candidate.get("title")),
        compact_value(candidate.get("description")),
    ]).casefold()
    if not text:
        return False
    query_l = query.casefold()
    if "formwork" in query_l:
        return "formwork" in text or "concrete form" in text or "concrete forms" in text
    if "concrete" in query_l:
        return "concrete" in text or "reinforced" in text
    if "tunnel" in query_l:
        return "tunnel" in text
    return any(token.casefold() in text for token in re.findall(r"[A-Za-z]{4,}", query)[:4])


def commons_image_candidates(query: str, limit: int = 8) -> list[dict[str, Any]]:
    params = {
        "action": "query",
        "generator": "search",
        "gsrsearch": query,
        "gsrnamespace": 6,
        "gsrlimit": str(limit),
        "prop": "imageinfo",
        "iiprop": "url|extmetadata|mime",
        "iiurlwidth": "1100",
        "format": "json",
    }
    try:
        response = requests.get(
            COMMONS_API_URL,
            params=params,
            timeout=25,
            headers={"User-Agent": COMMONS_USER_AGENT},
        )
        response.raise_for_status()
    except requests.RequestException:
        return []
    pages = response.json().get("query", {}).get("pages", {})
    candidates: list[dict[str, Any]] = []
    for page in pages.values():
        imageinfo = (page.get("imageinfo") or [{}])[0]
        mime = compact_value(imageinfo.get("mime"))
        if mime not in {"image/jpeg", "image/png"}:
            continue
        metadata = imageinfo.get("extmetadata") or {}
        license_label = strip_html((metadata.get("LicenseShortName") or {}).get("value"))
        usage_terms = strip_html((metadata.get("UsageTerms") or {}).get("value"))
        if not license_label and not usage_terms:
            continue
        image_url = imageinfo.get("thumburl") or imageinfo.get("url")
        if not image_url:
            continue
        candidates.append({
            "title": page.get("title", ""),
            "url": image_url,
            "description_url": imageinfo.get("descriptionurl"),
            "mime": mime,
            "license": license_label or usage_terms,
            "author": strip_html((metadata.get("Artist") or {}).get("value")),
            "description": strip_html((metadata.get("ImageDescription") or {}).get("value")),
        })
    return [candidate for candidate in candidates if relevant_commons_candidate(candidate, query)]


def image_dimensions(data: bytes, mime: str) -> tuple[int, int] | None:
    if mime == "image/png" and data.startswith(b"\x89PNG\r\n\x1a\n") and len(data) >= 24:
        return int.from_bytes(data[16:20], "big"), int.from_bytes(data[20:24], "big")

    if mime == "image/jpeg" and data.startswith(b"\xff\xd8"):
        index = 2
        while index + 9 < len(data):
            if data[index] != 0xFF:
                index += 1
                continue
            marker = data[index + 1]
            index += 2
            if marker in {0xD8, 0xD9}:
                continue
            if index + 2 > len(data):
                break
            segment_length = int.from_bytes(data[index:index + 2], "big")
            if segment_length < 2 or index + segment_length > len(data):
                break
            if 0xC0 <= marker <= 0xC3 and segment_length >= 7:
                height = int.from_bytes(data[index + 3:index + 5], "big")
                width = int.from_bytes(data[index + 5:index + 7], "big")
                return width, height
            index += segment_length

    return None


def valid_downloaded_image(data: bytes, mime: str) -> bool:
    dimensions = image_dimensions(data, mime)
    if not dimensions:
        return False
    width, height = dimensions
    return width >= 320 and height >= 180 and len(data) >= 4096


def download_commons_images(
    order: dict[str, Any],
    count: int,
    *,
    headings: list[str],
    workspace: Path | None,
) -> list[dict[str, Any]]:
    if count <= 0 or not workspace:
        return []
    figures_dir = workspace / "final" / "figures"
    figures_dir.mkdir(parents=True, exist_ok=True)

    selected: list[dict[str, Any]] = []
    seen_urls: set[str] = set()
    for query in image_search_queries(order, headings):
        for candidate in commons_image_candidates(query):
            if len(selected) >= count:
                break
            url = compact_value(candidate.get("url"))
            if not url or url in seen_urls:
                continue
            seen_urls.add(url)
            extension = ".png" if candidate.get("mime") == "image/png" else ".jpg"
            try:
                response = requests.get(
                    url,
                    timeout=40,
                    headers={"User-Agent": COMMONS_USER_AGENT},
                )
                response.raise_for_status()
            except requests.RequestException:
                continue
            if not valid_downloaded_image(response.content, compact_value(candidate.get("mime"))):
                continue
            index = len(selected) + 1
            file_path = figures_dir / f"web_image_{index}{extension}"
            file_path.write_bytes(response.content)
            source = {
                **candidate,
                "local_path": str(file_path.relative_to(workspace)),
                "query": query,
            }
            selected.append(source)
        if len(selected) >= count:
            break

    if selected:
        (figures_dir / "image_sources.json").write_text(
            json.dumps(selected, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
    return selected


def clean_generated_figures(workspace: Path) -> None:
    figures_dir = workspace / "final" / "figures"
    if not figures_dir.exists():
        return
    for pattern in ["auto_figure_*.*", "web_image_*.*"]:
        for path in figures_dir.glob(pattern):
            if path.is_file():
                path.unlink()


def resolve_image_path(raw_path: str, base_dir: Path | None, workspace: Path | None) -> Path | None:
    if re.match(r"^https?://", raw_path, flags=re.I):
        return None
    candidate = Path(raw_path)
    candidates = [candidate] if candidate.is_absolute() else []
    if not candidate.is_absolute():
        if base_dir:
            candidates.append(base_dir / candidate)
        if workspace:
            candidates.append(workspace / candidate)
    for item in candidates:
        if item.exists() and item.is_file():
            return item
    return None


def safe_svg_text(value: Any, limit: int = 56) -> str:
    text = compact_value(value)
    text = re.sub(r"\s+", " ", text)
    if len(text) > limit:
        text = text[: limit - 1].rstrip() + "…"
    return html.escape(text)


def short_plain_text(value: Any, limit: int = 76) -> str:
    text = compact_value(value)
    text = re.sub(r"\s+", " ", text)
    if len(text) > limit:
        text = text[: limit - 1].rstrip() + "…"
    return text


def web_image_caption(index: int, source: dict[str, Any]) -> str:
    query = compact_value(source.get("query")).casefold()
    if "formwork" in query and "tunnel" in query:
        topic = "نمایی از قالب‌بندی و آرماتوربندی بتن در پروژه تونلی"
    elif "formwork" in query:
        topic = "نمایی از قالب‌بندی بتن در پروژه ساختمانی"
    elif "concrete" in query:
        topic = "نمایی از اجرای بتن و اجزای سازه‌ای مرتبط"
    else:
        topic = "تصویر مرتبط با موضوع تحقیق"
    return f"شکل {index}: {topic}"


def markdown_blocks(markdown: str) -> list[tuple[str, Any]]:
    blocks: list[tuple[str, Any]] = []
    lines = markdown.splitlines()
    index = 0
    while index < len(lines):
        line = lines[index].strip()
        if not line:
            index += 1
            continue

        image = re.match(r"^!\[(.*?)\]\((.*?)\)\s*$", line)
        if image:
            blocks.append(("image", {"alt": image.group(1).strip(), "path": image.group(2).strip()}))
            index += 1
            continue

        if line.startswith("|"):
            table_lines = []
            while index < len(lines) and lines[index].strip().startswith("|"):
                table_lines.append(lines[index].strip())
                index += 1
            rows = []
            for table_line in table_lines:
                cells = [cell.strip() for cell in table_line.strip("|").split("|")]
                if cells and all(re.fullmatch(r":?-{3,}:?", cell or "") for cell in cells):
                    continue
                rows.append(cells)
            if rows:
                blocks.append(("table", rows))
            continue

        heading = re.match(r"^(#{1,6})\s+(.+)$", line)
        if heading:
            blocks.append(("heading", (len(heading.group(1)), heading.group(2).strip())))
            index += 1
            continue

        if re.match(r"^[-*]\s+", line):
            items = []
            while index < len(lines) and re.match(r"^[-*]\s+", lines[index].strip()):
                items.append(re.sub(r"^[-*]\s+", "", lines[index].strip()))
                index += 1
            blocks.append(("ul", items))
            continue

        if re.match(r"^\d+\.\s+", line):
            items = []
            while index < len(lines) and re.match(r"^\d+\.\s+", lines[index].strip()):
                items.append(re.sub(r"^\d+\.\s+", "", lines[index].strip()))
                index += 1
            blocks.append(("ol", items))
            continue

        paragraph = [line]
        index += 1
        while index < len(lines):
            next_line = lines[index].strip()
            if not next_line:
                index += 1
                break
            if next_line.startswith("|") or next_line.startswith("#") or re.match(r"^[-*]\s+", next_line) or re.match(r"^\d+\.\s+", next_line):
                break
            paragraph.append(next_line)
            index += 1
        blocks.append(("p", " ".join(paragraph)))
    return blocks


def block_weight(block: tuple[str, Any]) -> int:
    kind, value = block
    if kind == "heading":
        return 180
    if kind in {"image", "figure"}:
        return 700
    if kind == "table":
        return sum(sum(len(str(cell)) for cell in row) for row in value) + 600
    if kind in {"ul", "ol"}:
        return sum(len(item) for item in value) + 160
    return len(str(value))


def paginate_blocks(blocks: list[tuple[str, Any]], page_count: int | None) -> list[list[tuple[str, Any]]]:
    if not page_count or page_count <= 1:
        return [blocks]

    total_weight = max(sum(block_weight(block) for block in blocks), 1)
    target = max(total_weight // page_count, 900)
    pages: list[list[tuple[str, Any]]] = []
    current: list[tuple[str, Any]] = []
    current_weight = 0

    for block in blocks:
        remaining_blocks = len(blocks) - sum(len(page) for page in pages) - len(current)
        remaining_pages = page_count - len(pages)
        weight = block_weight(block)
        starts_major_section = block[0] == "heading" and block[1][0] <= 2 and current
        should_break = (
            len(pages) < page_count - 1
            and current
            and (current_weight + weight > target or (starts_major_section and current_weight > target * 0.55))
            and remaining_blocks >= remaining_pages
        )
        if should_break:
            pages.append(current)
            current = []
            current_weight = 0
        current.append(block)
        current_weight += weight

    if current:
        pages.append(current)
    while len(pages) < page_count:
        pages.append([("p", "")])
    return pages


def html_inline(value: str) -> str:
    escaped = html.escape(value)
    escaped = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", escaped)
    escaped = re.sub(r"\*(.+?)\*", r"<em>\1</em>", escaped)
    escaped = re.sub(r"`(.+?)`", r"<span class=\"ltr\">\1</span>", escaped)
    return escaped


def block_to_html(block: tuple[str, Any]) -> str:
    kind, value = block
    if kind == "heading":
        level, text = value
        tag = "h1" if level == 1 else "h2" if level == 2 else "h3"
        return f"<{tag}>{html_inline(text)}</{tag}>"
    if kind == "table":
        rows = value
        output = ["<table>"]
        for row_index, row in enumerate(rows):
            tag = "th" if row_index == 0 else "td"
            output.append("<tr>" + "".join(f"<{tag}>{html_inline(str(cell))}</{tag}>" for cell in row) + "</tr>")
        output.append("</table>")
        return "".join(output)
    if kind == "ul":
        return "<ul>" + "".join(f"<li>{html_inline(item)}</li>" for item in value) + "</ul>"
    if kind == "ol":
        return "<ol>" + "".join(f"<li>{html_inline(item)}</li>" for item in value) + "</ol>"
    if not str(value).strip():
        return "<p>&nbsp;</p>"
    return f"<p>{html_inline(str(value))}</p>"



PERSIAN_DIGITS = str.maketrans("0123456789", "۰۱۲۳۴۵۶۷۸۹")
RTL_MARK = "\u200f"


def persianize_numbering(text: str) -> str:
    # Leading Latin section/list numbers make Word visually place headings on the left
    # inside RTL paragraphs. Convert only numbering-like prefixes, not URLs/DOIs.
    def repl(match: re.Match[str]) -> str:
        return match.group(1).translate(PERSIAN_DIGITS) + match.group(2)

    text = re.sub(r"^(\d+(?:\.\d+)*)([\).\s-]+)", repl, text)
    text = re.sub(r"(?<=\s)(\d+(?:\.\d+)*)(?=\s*(?:تا|صفحه|روز|درصد|تکرار))", lambda m: m.group(1).translate(PERSIAN_DIGITS), text)
    return text


def rtl_display_text(text: str, rtl: bool) -> str:
    if rtl and contains_rtl(text):
        return RTL_MARK + persianize_numbering(text)
    return text


def word_run(text: str, *, bold: bool = False, size: int = 24, rtl: bool | None = None) -> str:
    if rtl is None:
        rtl = True
    escaped = html.escape(rtl_display_text(text, rtl))
    bold_xml = "<w:b/><w:bCs/>" if bold else ""
    rtl_xml = "<w:rtl/>" if rtl else ""
    return (
        "<w:r><w:rPr>"
        '<w:rFonts w:ascii="Times New Roman" w:hAnsi="Times New Roman" w:cs="B Nazanin"/>'
        f"{bold_xml}{rtl_xml}<w:sz w:val=\"{size}\"/><w:szCs w:val=\"{size}\"/>"
        "</w:rPr>"
        f"<w:t xml:space=\"preserve\">{escaped}</w:t></w:r>"
    )


def word_paragraph(
    text: str,
    *,
    heading_level: int | None = None,
    align: str | None = None,
) -> str:
    rtl = True
    if heading_level == 1:
        size = 32
        bold = True
        before = 120
        after = 160
        style = "Title"
        default_jc = "right"
    elif heading_level == 2:
        size = 28
        bold = True
        before = 120
        after = 60
        style = "Heading1"
        default_jc = "right"
    elif heading_level == 3:
        size = 26
        bold = True
        before = 100
        after = 50
        style = "Heading1"
        default_jc = "right"
    else:
        size = 22
        bold = False
        before = 0
        after = 70
        style = "Normal"
        default_jc = "right"

    jc = align or default_jc
    bidi = "<w:bidi/>"
    return (
        "<w:p><w:pPr>"
        f'<w:pStyle w:val="{style}"/>{bidi}<w:jc w:val="{jc}"/>'
        f'<w:spacing w:before="{before}" w:after="{after}" w:line="290" w:lineRule="auto"/>'
        "</w:pPr>"
        + word_run(text, bold=bold, size=size, rtl=rtl)
        + "</w:p>"
    )


def word_figure(asset: FigureAsset) -> str:
    cx = 5_120_000
    cy = 2_250_000
    doc_pr_id = re.sub(r"\D", "", asset.rid) or "1"
    drawing = (
        '<w:p><w:pPr><w:bidi/><w:jc w:val="center"/>'
        '<w:spacing w:before="90" w:after="50" w:line="260" w:lineRule="auto"/>'
        "</w:pPr><w:r><w:drawing>"
        '<wp:inline distT="0" distB="0" distL="0" distR="0">'
        f'<wp:extent cx="{cx}" cy="{cy}"/>'
        f'<wp:docPr id="{doc_pr_id}" name="{html.escape(asset.alt or asset.caption)}"/>'
        '<wp:cNvGraphicFramePr><a:graphicFrameLocks noChangeAspect="1"/></wp:cNvGraphicFramePr>'
        '<a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/picture">'
        '<pic:pic><pic:nvPicPr><pic:cNvPr id="0" name="figure"/>'
        '<pic:cNvPicPr/></pic:nvPicPr><pic:blipFill>'
        f'<a:blip r:embed="{asset.rid}"/><a:stretch><a:fillRect/></a:stretch>'
        '</pic:blipFill><pic:spPr><a:xfrm><a:off x="0" y="0"/>'
        f'<a:ext cx="{cx}" cy="{cy}"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        '</pic:spPr></pic:pic></a:graphicData></a:graphic>'
        '</wp:inline></w:drawing></w:r></w:p>'
    )
    caption = word_paragraph(asset.caption, align="center")
    return drawing + caption


def word_page_break() -> str:
    return '<w:p><w:r><w:br w:type="page"/></w:r></w:p>'


def word_table(rows: list[list[str]]) -> str:
    output = [
        "<w:tbl>",
        "<w:tblPr>",
        '<w:tblW w:w="0" w:type="auto"/>',
        '<w:jc w:val="right"/>',
        "<w:bidiVisual/>",
        '<w:tblBorders><w:top w:val="single" w:sz="6" w:space="0" w:color="666666"/>'
        '<w:left w:val="single" w:sz="6" w:space="0" w:color="666666"/>'
        '<w:bottom w:val="single" w:sz="6" w:space="0" w:color="666666"/>'
        '<w:right w:val="single" w:sz="6" w:space="0" w:color="666666"/>'
        '<w:insideH w:val="single" w:sz="6" w:space="0" w:color="666666"/>'
        '<w:insideV w:val="single" w:sz="6" w:space="0" w:color="666666"/></w:tblBorders>',
        "</w:tblPr>",
    ]
    for row_index, row in enumerate(rows):
        output.append("<w:tr>")
        for cell in row:
            shade = '<w:shd w:fill="EDEDED"/>' if row_index == 0 else ""
            output.append(
                "<w:tc><w:tcPr>"
                '<w:tcW w:w="2400" w:type="dxa"/>'
                f"{shade}"
                "</w:tcPr>"
                + word_paragraph(str(cell), heading_level=None, align="right")
                + "</w:tc>"
            )
        output.append("</w:tr>")
    output.append("</w:tbl>")
    return "".join(output)


def word_block(block: tuple[str, Any]) -> str:
    kind, value = block
    if kind == "heading":
        level, text = value
        return word_paragraph(text, heading_level=min(level, 3))
    if kind == "figure":
        return word_figure(value)
    if kind == "image":
        alt = value.get("alt") or "تصویر"
        return word_paragraph(f"[تصویر در دسترس نیست: {alt}]", align="center")
    if kind == "table":
        return word_table(value)
    if kind in {"ul", "ol"}:
        output = []
        for index, item in enumerate(value, start=1):
            prefix = f"{str(index).translate(PERSIAN_DIGITS)}. " if kind == "ol" else "• "
            output.append(word_paragraph(prefix + item, align="right"))
        return "".join(output)
    return word_paragraph(str(value)) if str(value).strip() else word_paragraph(" ")


def make_figure_asset(
    *,
    rid_index: int,
    data: bytes,
    extension: str,
    alt: str,
    caption: str,
    source_url: str | None = None,
    license_label: str | None = None,
    author: str | None = None,
) -> FigureAsset:
    extension = extension.lower()
    if extension == ".jpg":
        extension = ".jpeg"
    content_type = IMAGE_CONTENT_TYPES.get(extension, "image/svg+xml")
    target_extension = "jpg" if extension == ".jpeg" else extension.lstrip(".")
    return FigureAsset(
        rid=f"rIdImage{rid_index}",
        target=f"media/figure{rid_index}.{target_extension}",
        extension=target_extension,
        content_type=content_type,
        data=data,
        alt=alt,
        caption=caption,
        source_url=source_url,
        license_label=license_label,
        author=author,
    )


def prepare_docx_blocks(
    blocks: list[tuple[str, Any]],
    order: dict[str, Any],
    *,
    base_dir: Path | None,
    workspace: Path | None,
) -> tuple[list[tuple[str, Any]], list[FigureAsset]]:
    assets: list[FigureAsset] = []
    prepared: list[tuple[str, Any]] = []
    next_rid = 1

    for block in blocks:
        if block[0] != "image":
            prepared.append(block)
            continue

        value = block[1]
        raw_path = compact_value(value.get("path"))
        alt = compact_value(value.get("alt")) or "تصویر"
        image_path = resolve_image_path(raw_path, base_dir, workspace)
        if not image_path:
            prepared.append(block)
            continue
        extension = image_path.suffix.lower()
        if extension not in IMAGE_CONTENT_TYPES:
            prepared.append(("p", f"[فرمت تصویر پشتیبانی نمی‌شود: {image_path.name}]"))
            continue
        caption = alt if alt.startswith("شکل") else f"شکل {len(assets) + 1}: {alt}"
        asset = make_figure_asset(
            rid_index=next_rid,
            data=image_path.read_bytes(),
            extension=extension,
            alt=alt,
            caption=caption,
        )
        next_rid += 1
        assets.append(asset)
        prepared.append(("figure", asset))

    target_count = requested_or_default_image_count(order)
    if target_count <= len(assets):
        return prepared, assets

    heading_candidates = [
        value[1]
        for kind, value in prepared
        if kind == "heading"
        and value[0] >= 2
        and not any(skip in value[1] for skip in ["منابع", "چکیده", "واژگان", "فهرست"])
    ]
    if not heading_candidates:
        heading_candidates = [compact_value(order.get("title")) or "موضوع اصلی"]

    insert_after_indexes: dict[int, list[FigureAsset]] = {}
    heading_positions = [i for i, block in enumerate(prepared) if block[0] == "heading" and block[1][1] in heading_candidates]
    if not heading_positions:
        heading_positions = [min(len(prepared) - 1, 0)]

    web_sources = download_commons_images(
        order,
        target_count - len(assets),
        headings=heading_candidates,
        workspace=workspace,
    )
    for source in web_sources:
        image_path = workspace / source["local_path"] if workspace else None
        if not image_path or not image_path.exists():
            continue
        extension = image_path.suffix.lower()
        caption_source = compact_value(source.get("title")).removeprefix("File:").strip()
        caption = web_image_caption(len(assets) + 1, source)
        asset = make_figure_asset(
            rid_index=next_rid,
            data=image_path.read_bytes(),
            extension=extension,
            alt=caption_source or "تصویر مرتبط با موضوع",
            caption=caption,
            source_url=source.get("description_url"),
            license_label=source.get("license"),
            author=source.get("author"),
        )
        next_rid += 1
        assets.append(asset)
        position = heading_positions[(len(assets) - 1) % len(heading_positions)]
        insert_after_indexes.setdefault(position, []).append(asset)

    output: list[tuple[str, Any]] = []
    for index, block in enumerate(prepared):
        output.append(block)
        for asset in insert_after_indexes.get(index, []):
            output.append(("figure", asset))
    return output, assets


def word_styles_xml() -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:styles xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        '<w:latentStyles w:defLockedState="0" w:defUIPriority="99" w:defSemiHidden="0" '
        'w:defUnhideWhenUsed="0" w:defQFormat="0" w:count="276"/>'
        '<w:docDefaults>'
        '<w:rPrDefault><w:rPr>'
        '<w:rFonts w:ascii="Times New Roman" w:hAnsi="Times New Roman" w:cs="B Nazanin"/>'
        '<w:sz w:val="22"/><w:szCs w:val="22"/><w:lang w:val="en-US" w:bidi="fa-IR"/>'
        '</w:rPr></w:rPrDefault>'
        '<w:pPrDefault><w:pPr><w:bidi/><w:jc w:val="right"/>'
        '<w:spacing w:after="70" w:line="290" w:lineRule="auto"/>'
        '</w:pPr></w:pPrDefault>'
        '</w:docDefaults>'
        '<w:style w:type="paragraph" w:default="1" w:styleId="Normal">'
        '<w:name w:val="Normal"/><w:qFormat/>'
        '<w:pPr><w:bidi/><w:jc w:val="right"/><w:spacing w:after="70" w:line="290" w:lineRule="auto"/></w:pPr>'
        '<w:rPr><w:rFonts w:ascii="Times New Roman" w:hAnsi="Times New Roman" w:cs="B Nazanin"/>'
        '<w:rtl/><w:sz w:val="22"/><w:szCs w:val="22"/><w:lang w:bidi="fa-IR"/></w:rPr>'
        '</w:style>'
        '<w:style w:type="paragraph" w:styleId="Heading1">'
        '<w:name w:val="heading 1"/><w:basedOn w:val="Normal"/><w:next w:val="Normal"/><w:qFormat/>'
        '<w:pPr><w:bidi/><w:jc w:val="right"/><w:keepNext/><w:spacing w:before="120" w:after="60" w:line="290" w:lineRule="auto"/></w:pPr>'
        '<w:rPr><w:rFonts w:ascii="Times New Roman" w:hAnsi="Times New Roman" w:cs="B Nazanin"/>'
        '<w:b/><w:bCs/><w:rtl/><w:sz w:val="28"/><w:szCs w:val="28"/><w:lang w:bidi="fa-IR"/></w:rPr>'
        '</w:style>'
        '<w:style w:type="paragraph" w:styleId="Title">'
        '<w:name w:val="Title"/><w:basedOn w:val="Normal"/><w:next w:val="Normal"/><w:qFormat/>'
        '<w:pPr><w:bidi/><w:jc w:val="right"/><w:keepNext/><w:spacing w:before="120" w:after="160" w:line="290" w:lineRule="auto"/></w:pPr>'
        '<w:rPr><w:rFonts w:ascii="Times New Roman" w:hAnsi="Times New Roman" w:cs="B Nazanin"/>'
        '<w:b/><w:bCs/><w:rtl/><w:sz w:val="32"/><w:szCs w:val="32"/><w:lang w:bidi="fa-IR"/></w:rPr>'
        '</w:style>'
        '</w:styles>'
    )


def word_settings_xml() -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:settings xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        '<w:bidi/>'
        '<w:defaultTabStop w:val="720"/>'
        '<w:themeFontLang w:val="en-US" w:eastAsia="en-US" w:bidi="fa-IR"/>'
        '<w:decimalSymbol w:val="."/><w:listSeparator w:val=","/>'
        '<w:compat><w:compatSetting w:name="compatibilityMode" '
        'w:uri="http://schemas.microsoft.com/office/word" w:val="15"/></w:compat>'
        '</w:settings>'
    )


def write_native_formatted_docx(
    path: Path,
    title: str,
    markdown: str,
    order: dict[str, Any],
    *,
    base_dir: Path | None = None,
    workspace: Path | None = None,
) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    # Real deliverables must flow naturally. Page count is controlled by source
    # length and style, not by inserting artificial page breaks that leave gaps.
    blocks, image_assets = prepare_docx_blocks(
        markdown_blocks(markdown),
        order,
        base_dir=base_dir,
        workspace=workspace,
    )
    body_parts = [word_block(block) for block in blocks]

    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:pic="http://schemas.openxmlformats.org/drawingml/2006/picture">'
        "<w:body>"
        + "".join(body_parts)
        + '<w:sectPr><w:pgSz w:w="11906" w:h="16838"/>'
        '<w:pgMar w:top="900" w:right="900" w:bottom="900" w:left="900" w:header="720" w:footer="720" w:gutter="0"/>'
        '<w:bidi/><w:rtlGutter/><w:cols w:space="720"/><w:docGrid w:linePitch="360"/></w:sectPr>'
        + "</w:body></w:document>"
    )
    image_defaults = {
        asset.extension: asset.content_type
        for asset in image_assets
    }
    content_types = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        + "".join(
            f'<Default Extension="{extension}" ContentType="{content_type}"/>'
            for extension, content_type in sorted(image_defaults.items())
        )
        + '<Override PartName="/word/document.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        '<Override PartName="/word/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.styles+xml"/>'
        '<Override PartName="/word/settings.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.settings+xml"/>'
        "</Types>"
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
        'Target="word/document.xml"/>'
        "</Relationships>"
    )
    document_rels = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rIdStyles" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles" Target="styles.xml"/>'
        '<Relationship Id="rIdSettings" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/settings" Target="settings.xml"/>'
        + "".join(
            f'<Relationship Id="{asset.rid}" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" '
            f'Target="{asset.target}"/>'
            for asset in image_assets
        )
        + '</Relationships>'
    )
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as docx:
        docx.writestr("[Content_Types].xml", content_types)
        docx.writestr("_rels/.rels", rels)
        docx.writestr("word/_rels/document.xml.rels", document_rels)
        docx.writestr("word/document.xml", document_xml)
        docx.writestr("word/styles.xml", word_styles_xml())
        docx.writestr("word/settings.xml", word_settings_xml())
        for asset in image_assets:
            docx.writestr(f"word/{asset.target}", asset.data)


def write_rtl_html_docx(path: Path, title: str, markdown: str, order: dict[str, Any]) -> bool:
    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not soffice:
        return False

    page_count = requested_page_count(order)
    blocks = markdown_blocks(markdown)
    pages = paginate_blocks(blocks, page_count)
    page_html = "\n".join(
        '<section class="page">' + "\n".join(block_to_html(block) for block in page) + "</section>"
        for page in pages
    )
    html_text = f"""<!doctype html>
<html lang="fa" dir="rtl">
<head>
  <meta charset="utf-8">
  <title>{html.escape(title)}</title>
  <style>
    @page {{ size: A4; margin: 2.2cm 2.0cm 2.2cm 2.0cm; }}
    body {{
      direction: rtl;
      unicode-bidi: isolate;
      font-family: "B Nazanin", "Vazirmatn", "IRANSans", Tahoma, Arial, sans-serif;
      font-size: 12.5pt;
      line-height: 1.45;
      color: #111;
      text-align: right;
    }}
    .page {{ page-break-after: always; }}
    .page:last-child {{ page-break-after: auto; }}
    h1 {{
      font-size: 17pt;
      text-align: center;
      margin: 0 0 14pt;
      font-weight: 700;
    }}
    h2 {{
      font-size: 14pt;
      margin: 12pt 0 6pt;
      font-weight: 700;
    }}
    h3 {{
      font-size: 13pt;
      margin: 10pt 0 5pt;
      font-weight: 700;
    }}
    p {{ margin: 0 0 7pt; text-align: justify; }}
    table {{
      width: 100%;
      border-collapse: collapse;
      margin: 8pt 0 10pt;
      direction: rtl;
      font-size: 10.5pt;
    }}
    th, td {{
      border: 1px solid #555;
      padding: 4pt 5pt;
      vertical-align: top;
      text-align: right;
    }}
    th {{ background: #eaeaea; font-weight: 700; }}
    ul, ol {{ margin: 0 18pt 8pt 0; padding: 0; }}
    li {{ margin-bottom: 4pt; }}
    .ltr {{
      direction: ltr;
      unicode-bidi: isolate;
      font-family: "Times New Roman", Calibri, Arial, sans-serif;
    }}
  </style>
</head>
<body>
{page_html}
</body>
</html>
"""
    path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="payanname-docx-") as tmpdir:
        tmp = Path(tmpdir)
        html_path = tmp / "deliverable.html"
        html_path.write_text(html_text, encoding="utf-8")
        subprocess.run(
            [soffice, "--headless", "--convert-to", "docx", "--outdir", str(tmp), str(html_path)],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )
        converted = tmp / "deliverable.docx"
        if not converted.exists():
            return False
        shutil.copyfile(converted, path)
    return True


def slide_sections_from_markdown(markdown: str, max_slides: int | None = None) -> list[dict[str, Any]]:
    sections: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    in_references = False

    for raw_line in markdown.splitlines():
        line = raw_line.strip()
        if not line or line.startswith("<"):
            continue

        heading = re.match(r"^###\s+(.+)$", line)
        if heading:
            title = re.sub(r"^\d+[\.\)]\s*", "", heading.group(1).strip())
            current = {"title": title, "bullets": [], "notes": []}
            sections.append(current)
            in_references = "منابع" in title
            continue

        if current is None or in_references:
            continue

        if line.startswith("منبع:") or line.startswith("شکل "):
            current["notes"].append(line)
            continue

        if line.startswith("|"):
            cells = [cell.strip() for cell in line.strip("|").split("|")]
            if cells and all(re.fullmatch(r":?-{3,}:?", cell or "") for cell in cells):
                continue
            if len(cells) >= 4 and cells[0] != "معیار":
                current["bullets"].append(
                    f"{cells[0]}: قالب تونلی {cells[1]}؛ قالب آلومینیومی {cells[2]}؛ قالب سنتی {cells[3]}"
                )
            continue

        if re.match(r"^[-*]\s+", line):
            bullet = re.sub(r"^[-*]\s+", "", line)
            if bullet:
                current["bullets"].append(bullet)

    if max_slides:
        sections = sections[:max_slides]
    return sections


def load_sourced_image_paths(workspace: Path) -> list[Path]:
    sources_path = workspace / "final" / "figures" / "image_sources.json"
    if not sources_path.exists():
        return []
    try:
        sources = json.loads(sources_path.read_text(encoding="utf-8"))
    except (OSError, ValueError):
        return []
    paths: list[Path] = []
    if not isinstance(sources, list):
        return paths
    for source in sources:
        if not isinstance(source, dict):
            continue
        if compact_value(source.get("generated")) or compact_value(source.get("generation_tool")):
            continue
        local_path = compact_value(source.get("local_path"))
        candidate = workspace / local_path if local_path else None
        if candidate and candidate.exists() and candidate.suffix.lower() in {".jpg", ".jpeg", ".png"}:
            paths.append(candidate)
    return paths


def load_sourced_figures(workspace: Path) -> list[dict[str, Any]]:
    sources_path = workspace / "final" / "figures" / "image_sources.json"
    if not sources_path.exists():
        return []
    try:
        sources = json.loads(sources_path.read_text(encoding="utf-8"))
    except (OSError, ValueError):
        return []
    figures: list[dict[str, Any]] = []
    if not isinstance(sources, list):
        return figures
    for source in sources:
        if not isinstance(source, dict):
            continue
        if compact_value(source.get("generated")) or compact_value(source.get("generation_tool")):
            continue
        local_path = compact_value(source.get("local_path"))
        candidate = workspace / local_path if local_path else None
        if candidate and candidate.exists() and candidate.suffix.lower() in {".jpg", ".jpeg", ".png"}:
            figures.append({**source, "path": candidate})
    return figures


def slide_image_for(index: int, section: dict[str, Any], figures: list[dict[str, Any]]) -> dict[str, Any] | None:
    if not figures:
        return None
    title = compact_value(section.get("title"))
    local_paths = [compact_value(item.get("local_path")) for item in figures]

    def find(pattern: str) -> dict[str, Any] | None:
        for figure, local_path in zip(figures, local_paths):
            if pattern in local_path:
                return figure
        return None

    if index == 2:
        return find("photo") or figures[0]
    if index == 3:
        return find("components") or figures[min(1, len(figures) - 1)]
    if index == 4:
        return find("cycle_bar_chart_readable") or find("cycle") or figures[min(2, len(figures) - 1)]
    if "اجزا" in title:
        return find("components")
    if "چرخه" in title:
        return find("cycle_bar_chart_readable") or find("cycle")
    return None


def title_slide_figure(figures: list[dict[str, Any]]) -> dict[str, Any] | None:
    for figure in figures:
        local_path = compact_value(figure.get("local_path"))
        if "photo" in local_path:
            return figure
    return figures[0] if figures else None


def fit_picture(slide: Any, image_path: Path, left: Any, top: Any, width: Any, height: Any) -> Any:
    from pptx.util import Inches
    from PIL import Image

    with Image.open(image_path) as image:
        image_width, image_height = image.size
    if image_width <= 0 or image_height <= 0:
        return slide.shapes.add_picture(str(image_path), left, top, width=width)

    box_w = int(width)
    box_h = int(height)
    ratio = min(box_w / image_width, box_h / image_height)
    final_w = int(image_width * ratio)
    final_h = int(image_height * ratio)
    final_left = int(left) + max((box_w - final_w) // 2, 0)
    final_top = int(top) + max((box_h - final_h) // 2, 0)
    return slide.shapes.add_picture(str(image_path), final_left, final_top, width=final_w, height=final_h)


def set_shape_fill(shape: Any, rgb: tuple[int, int, int]) -> None:
    from pptx.dml.color import RGBColor

    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)
    shape.line.fill.background()


def add_textbox(
    slide: Any,
    text: str,
    *,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    font_size: int,
    bold: bool = False,
    color: tuple[int, int, int] = (28, 36, 46),
    align: Any = None,
    font_name: str = "B Nazanin",
) -> Any:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align or PP_ALIGN.RIGHT
    run = paragraph.add_run()
    run.text = rtl_display_text(text, True)
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_bullets(slide: Any, bullets: list[str], *, left: Any, top: Any, width: Any, height: Any) -> Any:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_right = Inches(0.12)
    frame.margin_left = Inches(0.08)
    for bullet_index, bullet in enumerate(bullets[:5]):
        paragraph = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.RIGHT
        paragraph.level = 0
        run = paragraph.add_run()
        run.text = rtl_display_text(f"• {bullet}", True)
        run.font.name = "B Nazanin"
        run.font.size = Pt(17 if len(bullets) <= 4 else 15)
        run.font.color.rgb = RGBColor(31, 42, 55)
    return box


def write_presentation_from_markdown(path: Path, title: str, markdown: str, order: dict[str, Any], workspace: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise SystemExit("python-pptx is required to generate presentation deliverables. Run `pip install -r requirements.txt`.") from exc

    try:
        max_slides = int(order.get("quantity_value") or 0) if order.get("quantity_type") == "slides" else None
    except (TypeError, ValueError):
        max_slides = None

    sections = slide_sections_from_markdown(markdown, max_slides=max_slides)
    if not sections:
        sections = [{"title": title, "bullets": markdown_to_plain_lines(markdown)[:5], "notes": []}]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    figures = load_sourced_figures(workspace)
    palette = {
        "ink": (29, 37, 47),
        "muted": (92, 103, 115),
        "green": (35, 116, 100),
        "amber": (207, 134, 45),
        "panel": (247, 249, 248),
        "line": (218, 224, 221),
    }

    for index, section in enumerate(sections, start=1):
        slide = prs.slides.add_slide(blank)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        set_shape_fill(bg, (255, 255, 255))
        bg.z_order = 0

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.82), 0, Inches(0.18), prs.slide_height)
        set_shape_fill(accent, palette["green"])

        bullets = list(section.get("bullets") or [])[:5]
        if not bullets:
            bullets = [
                compact_value(order.get("title")) or "موضوع ارائه",
                "جمع‌بندی آموزشی بر اساس منابع ثبت‌شده",
                "بدون عددسازی یا ادعای فاقد منبع",
            ]

        if index == 1:
            hero = title_slide_figure(figures)
            visual_panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.55), Inches(5.55), Inches(6.25))
            set_shape_fill(visual_panel, palette["panel"])
            if hero:
                fit_picture(slide, Path(hero["path"]), Inches(0.72), Inches(0.78), Inches(5.22), Inches(5.75))
            title_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.35), Inches(1.15), Inches(5.95), Inches(4.95))
            set_shape_fill(title_panel, palette["green"])
            add_textbox(
                slide,
                str(section["title"]),
                left=Inches(6.72),
                top=Inches(2.0),
                width=Inches(5.2),
                height=Inches(0.82),
                font_size=Pt(34),
                bold=True,
                color=(255, 255, 255),
                align=PP_ALIGN.CENTER,
            )
            subtitle = compact_value(order.get("course_name")) or compact_value(order.get("order_type")) or "ارائه"
            add_textbox(
                slide,
                subtitle,
                left=Inches(6.72),
                top=Inches(3.02),
                width=Inches(5.2),
                height=Inches(0.42),
                font_size=Pt(16),
                color=(232, 244, 240),
                align=PP_ALIGN.CENTER,
            )
            meta = " | ".join(
                value for value in [
                    compact_value(order.get("university")),
                    compact_value(order.get("field_of_study")),
                    compact_value(order.get("instructor_name")),
                ]
                if value
            )
            add_textbox(
                slide,
                meta,
                left=Inches(6.72),
                top=Inches(3.62),
                width=Inches(5.2),
                height=Inches(0.7),
                font_size=Pt(13),
                color=(232, 244, 240),
                align=PP_ALIGN.CENTER,
            )
            chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(6.86), Inches(1.28), Inches(0.32))
            set_shape_fill(chip, palette["green"])
            add_textbox(
                slide,
                f"{index} / {len(sections)}",
                left=Inches(0.72),
                top=Inches(6.88),
                width=Inches(1.12),
                height=Inches(0.24),
                font_size=Pt(9),
                bold=True,
                color=(255, 255, 255),
                align=PP_ALIGN.CENTER,
            )
            continue

        top_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(0.92), Inches(12.0), Inches(0.03))
        set_shape_fill(top_rule, palette["line"])

        add_textbox(
            slide,
            str(section["title"]),
            left=Inches(0.65),
            top=Inches(0.26),
            width=Inches(11.65),
            height=Inches(0.56),
            font_size=Pt(25),
            bold=True,
            color=palette["ink"],
        )

        figure = slide_image_for(index, section, figures)
        if figure:
            is_cycle_chart = "cycle" in compact_value(figure.get("local_path"))
            image_left = Inches(0.6)
            image_top = Inches(1.18)
            image_width = Inches(6.25 if is_cycle_chart else 5.15)
            image_height = Inches(5.55)
            image_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, image_left, image_top, image_width, image_height)
            set_shape_fill(image_panel, palette["panel"])
            fit_picture(
                slide,
                Path(figure["path"]),
                image_left + Inches(0.22),
                image_top + Inches(0.28),
                image_width - Inches(0.44),
                Inches(4.6 if is_cycle_chart else 4.65),
            )
            source_label = compact_value(figure.get("caption")) or compact_value(figure.get("title")) or "تصویر منبع‌دار"
            add_textbox(
                slide,
                f"منبع: {short_plain_text(source_label, 62)}",
                left=image_left + Inches(0.22),
                top=Inches(6.08),
                width=image_width - Inches(0.44),
                height=Inches(0.42),
                font_size=Pt(9),
                color=palette["muted"],
                align=PP_ALIGN.CENTER,
            )
            content_left = Inches(7.08 if is_cycle_chart else 6.05)
            content_width = Inches(5.52 if is_cycle_chart else 6.55)
        else:
            content_left = Inches(1.15)
            content_width = Inches(11.15)

        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, content_left, Inches(1.18), content_width, Inches(5.55))
        set_shape_fill(panel, palette["panel"])
        add_bullets(slide, bullets, left=content_left + Inches(0.35), top=Inches(1.55), width=content_width - Inches(0.7), height=Inches(4.55))

        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(6.86), Inches(1.28), Inches(0.32))
        set_shape_fill(chip, palette["green"])
        add_textbox(
            slide,
            f"{index} / {len(sections)}",
            left=Inches(0.72),
            top=Inches(6.88),
            width=Inches(1.12),
            height=Inches(0.24),
            font_size=Pt(9),
            bold=True,
            color=(255, 255, 255),
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            compact_value(order.get("course_name")) or compact_value(order.get("order_type")) or "ارائه",
            left=Inches(2.0),
            top=Inches(6.88),
            width=Inches(10.3),
            height=Inches(0.24),
            font_size=Pt(9),
            color=palette["muted"],
        )

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)


def package_existing(config: Config, args: argparse.Namespace) -> None:
    order_id, workspace = resolve_order_context(config, args.order_id, args.workspace)
    customer_input = workspace / "customer_input.json"
    order = json.loads(customer_input.read_text(encoding="utf-8")) if customer_input.exists() else {}
    title = order.get("title") or "Packaged Academic Deliverable"

    source_path = workspace / "final" / "deliverable_source.md"
    if args.source:
        source_path = normalize_workspace(args.source) or Path(args.source)
    elif not source_path.exists():
        legacy_source = workspace / "final" / "thesis_source.md"
        if legacy_source.exists():
            source_path = legacy_source
        else:
            source_path = workspace / "drafts" / "assisted_draft.md"

    if not source_path.exists():
        raise SystemExit(
            "No source found. Expected final/deliverable_source.md, final/thesis_source.md, or drafts/assisted_draft.md."
        )

    markdown = source_path.read_text(encoding="utf-8")
    docx_path = workspace / "final" / "deliverable.docx"
    clean_generated_figures(workspace)
    write_native_formatted_docx(
        docx_path,
        title,
        markdown,
        order,
        base_dir=source_path.parent,
        workspace=workspace,
    )
    validate_docx_rtl_quality(order, str(docx_path))
    generated = ["final/deliverable.docx"]
    if is_presentation_order(order):
        pptx_path = workspace / "final" / "deliverable.pptx"
        write_presentation_from_markdown(pptx_path, title, markdown, order, workspace)
        generated.append("final/deliverable.pptx")
    stale_pdf = workspace / "final" / "deliverable.pdf"
    if stale_pdf.exists():
        stale_pdf.unlink()

    print(
        json.dumps(
            {
                "orderId": order_id,
                "source": str(source_path),
                "generated": generated,
            },
            ensure_ascii=False,
            indent=2,
        )
    )


def current(config: Config, args: argparse.Namespace) -> None:
    order_id, workspace = resolve_order_context(config, args.order_id, args.workspace)
    print(
        json.dumps(
            {"orderId": order_id, "workspace": str(workspace), "workerId": config.worker_id},
            ensure_ascii=False,
            indent=2,
        )
    )


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Local Payanname worker helper")
    subparsers = parser.add_subparsers(dest="command", required=True)

    subparsers.add_parser(
        "run", help="Claim oldest approved order, start it, and create local workspace"
    )
    subparsers.add_parser(
        "claim", help="Alias for run: claim oldest approved order and create local workspace"
    )

    current_parser = subparsers.add_parser("current", help="Show order context for a workspace")
    current_parser.add_argument("--order-id")
    current_parser.add_argument("--workspace")

    heartbeat_parser = subparsers.add_parser("heartbeat", help="Refresh an active worker lock")
    heartbeat_parser.add_argument("--order-id")
    heartbeat_parser.add_argument("--workspace")

    draft_parser = subparsers.add_parser("submit-draft", help="Upload a draft file for review")
    draft_parser.add_argument("--order-id")
    draft_parser.add_argument("--workspace")
    draft_parser.add_argument("--draft-file")
    draft_parser.add_argument("--notes")

    final_parser = subparsers.add_parser("submit-final", help="Upload the review package for admin review")
    final_parser.add_argument("--order-id")
    final_parser.add_argument("--workspace")
    final_parser.add_argument("--pptx")
    final_parser.add_argument("--docx")
    final_parser.add_argument("--pdf")
    final_parser.add_argument("--deliverable-source")
    final_parser.add_argument("--compliance-report")
    final_parser.add_argument("--reference-usage-report")
    final_parser.add_argument("--human-review-checklist")
    final_parser.add_argument("--final-readme")
    final_parser.add_argument("--image-sources")
    final_parser.add_argument(
        "--skip-package-check",
        action="store_true",
        help="Submit whatever outputs are present without enforcing the standard review-package file set",
    )
    final_parser.add_argument("--notes")
    final_parser.add_argument(
        "--replace-existing",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Replace existing final_outputs of the same type for this order. Enabled by default.",
    )

    reset_parser = subparsers.add_parser("reset", help="Reset an interrupted order back to the approved pickup list")
    reset_parser.add_argument("--order-id")
    reset_parser.add_argument("--workspace")
    reset_parser.add_argument("--notes")

    mock_parser = subparsers.add_parser(
        "mock-generate", help="Generate a tiny mock academic package for system tests"
    )
    mock_parser.add_argument("--order-id")
    mock_parser.add_argument("--workspace")

    package_parser = subparsers.add_parser(
        "package-existing", help="Create editable DOCX from an existing Markdown academic source"
    )
    package_parser.add_argument("--order-id")
    package_parser.add_argument("--workspace")
    package_parser.add_argument("--source")

    fail_parser = subparsers.add_parser("fail", help="Mark order failed")
    fail_parser.add_argument("--order-id")
    fail_parser.add_argument("--workspace")
    fail_parser.add_argument("--notes")

    return parser


def main() -> None:
    config = load_config()
    parser = build_parser()
    args = parser.parse_args()

    config.workspace.mkdir(parents=True, exist_ok=True)

    try:
        if args.command in {"run", "claim"}:
            claim(config)
        elif args.command == "current":
            current(config, args)
        elif args.command == "heartbeat":
            heartbeat(config, args)
        elif args.command == "submit-draft":
            submit_draft(config, args)
        elif args.command == "submit-final":
            submit_final(config, args)
        elif args.command == "reset":
            reset(config, args)
        elif args.command == "mock-generate":
            mock_generate(config, args)
        elif args.command == "package-existing":
            package_existing(config, args)
        elif args.command == "fail":
            fail(config, args)
        else:
            parser.print_help()
            sys.exit(2)
    except requests.RequestException as exc:
        raise SystemExit(f"Worker request failed: {exc}") from exc


if __name__ == "__main__":
    main()
